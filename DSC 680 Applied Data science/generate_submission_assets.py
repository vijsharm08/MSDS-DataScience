from pathlib import Path
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import pyttsx3

ROOT = Path(r"c:\Users\us66099\Downloads\Bellevue University\DSC 680")
REPORTS = ROOT / "reports"
OUTPUT = ROOT / "output"
MODELS = ROOT / "models"

REPORT_DOCX = REPORTS / "DSC680_Final_Report.docx"
PRESENTATION_PPTX = REPORTS / "DSC680_Final_Presentation.pptx"
QA_DOCX = REPORTS / "DSC680_Q_and_A.docx"
AUDIO_WAV = REPORTS / "DSC680_Presentation_Audio.wav"
AUDIO_SCRIPT = REPORTS / "DSC680_Presentation_Script.txt"
CODE_NOTES = REPORTS / "CODE_SUBMISSION_NOTES.md"


def build_report_docx():
    doc = Document()
    styles = doc.styles
    styles['Normal'].font.name = 'Calibri'
    styles['Normal'].font.size = Pt(11)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run('Predicting High-Cost Healthcare Members')
    run.bold = True
    run.font.size = Pt(20)

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.add_run('DSC680-T301 Applied Data Science Final Project Report').italic = True

    doc.add_paragraph()
    p = doc.add_paragraph()
    p.add_run('Executive Summary').bold = True
    p = doc.add_paragraph()
    p.add_run('This white paper presents a data science workflow for identifying members likely to generate high healthcare costs. The project uses a realistic healthcare-style dataset, evaluates multiple regression models, and recommends a practical approach for targeted care intervention. The final model achieved strong predictive performance, with a Ridge Regression model producing an R² value of 0.9687 and an RMSE of $1,229.62, demonstrating strong accuracy for cost estimation.')

    doc.add_paragraph()
    p = doc.add_paragraph()
    p.add_run('Project Objective').bold = True
    p = doc.add_paragraph()
    p.add_run('The objective was to support proactive healthcare management by estimating future member cost and flagging members at elevated risk of high spending. This allows organizations to focus care-management resources more effectively and reduce avoidable costs.')

    doc.add_paragraph()
    p = doc.add_paragraph()
    p.add_run('Methodology').bold = True
    p = doc.add_paragraph()
    p.add_run('The workflow followed a standard supervised learning pipeline: data validation, exploratory analysis, feature engineering, model training, and evaluation. The project used fields such as age, gender, visit counts, emergency room visits, hospitalizations, and cost. Several engineered features were created, including age group, visits-per-month, and utilization interaction features.')

    doc.add_paragraph()
    p = doc.add_paragraph()
    p.add_run('Key Findings').bold = True
    p = doc.add_paragraph()
    p.add_run('The analysis showed that healthcare costs are strongly associated with utilization characteristics, especially hospitalizations and emergency room visits. A Pareto-style analysis also showed that a relatively small share of members accounted for a large share of total healthcare spend, supporting the value of targeted intervention strategies.')

    doc.add_paragraph()
    p = doc.add_paragraph()
    p.add_run('Results').bold = True
    p = doc.add_paragraph()
    p.add_run('The Ridge Regression model outperformed other tested approaches and delivered a strong fit on unseen data. The mean prediction error was very small, indicating that the model provides dependable estimates for practical decision support.')

    doc.add_paragraph()
    p = doc.add_paragraph()
    p.add_run('Ethical Considerations').bold = True
    p = doc.add_paragraph()
    p.add_run('Ethical use of this model requires careful attention to patient privacy, data governance, and transparency. Because the project uses a synthetic healthcare-style dataset rather than identifiable real patient records, it is appropriate for academic modeling, but any future deployment should use de-identified or properly consented real-world data and follow institutional review and compliance requirements. The model should also be monitored for bias across demographic groups, since healthcare utilization patterns can differ by age, gender, race, income, or geography. Fairness must be evaluated before the model is used to guide care decisions, and human clinicians should remain responsible for final decisions. Predictive models should support, not replace, clinical judgment, especially when interventions may affect access to care or resource allocation.')

    doc.add_paragraph()
    p = doc.add_paragraph()
    p.add_run('Conclusion').bold = True
    p = doc.add_paragraph()
    p.add_run('The project demonstrates how data science can support healthcare organizations in identifying high-cost members earlier and planning more targeted interventions. Although the dataset is synthetic, the workflow reflects a realistic and valuable analytics process for cost prediction and care management.')

    doc.add_paragraph()
    p = doc.add_paragraph()
    p.add_run('Artifacts').bold = True
    p = doc.add_paragraph()
    p.add_run('The project outputs include cost distribution and Pareto charts, a correlation heatmap, a validated prediction file, and a trained Ridge Regression model saved in the project folders.')

    doc.save(REPORT_DOCX)


def build_presentation_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle
        return slide

    def add_bullet_slide(title, bullets, image_path=None):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        body.clear()
        for i, bullet in enumerate(bullets):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(22)
            p.font.name = 'Calibri'
            p.alignment = PP_ALIGN.LEFT
        if image_path:
            slide.shapes.add_picture(str(image_path), Inches(8.2), Inches(1.3), width=Inches(4.4), height=Inches(3.0))
        return slide

    add_title_slide('Predicting High-Cost Healthcare Members', 'DSC680 Final Project Presentation')
    add_bullet_slide('Why this project matters', [
        'Healthcare organizations spend disproportionate amounts on a small subset of members.',
        'Earlier identification of high-cost members enables targeted interventions.',
        'The project converts healthcare utilization patterns into predictive insights.'
    ], OUTPUT / 'cost_distribution.png')
    add_bullet_slide('Data and methodology', [
        'Used a realistic public healthcare-style dataset with 1,000 member records.',
        'Explored age, gender, visits, ER use, hospitalizations, and cost.',
        'Applied feature engineering, model comparison, and evaluation metrics such as R² and RMSE.'
    ], OUTPUT / 'correlation_heatmap.png')
    add_bullet_slide('Key findings', [
        'Hospitalizations and emergency room visits were the strongest cost drivers.',
        'The distribution of costs was highly skewed and concentrated.',
        'A Pareto-style analysis confirmed that a subset of members produced a large share of spend.'
    ], OUTPUT / 'pareto_analysis.png')
    add_bullet_slide('Model performance', [
        'Ridge Regression achieved an R² of 0.9687 on the test set.',
        'RMSE was 1,229.62 and MAE was 983.96.',
        'Predicted costs closely aligned with the actual values, supporting practical use.'
    ])
    add_bullet_slide('Ethical considerations', [
        'Patient privacy and responsible data governance are essential in healthcare analytics.',
        'Synthetic data is useful for model development, but real-world deployment should use de-identified and compliant datasets.',
        'Models should be tested for bias across demographic groups and used to support fair, human-led care decisions.'
    ])
    add_bullet_slide('Business value and next steps', [
        'The model can support proactive care management and resource prioritization.',
        'Future work can include larger datasets, additional clinical features, and deployment into a dashboard.',
        'This project provides a strong foundation for continuing healthcare analytics work.'
    ])

    prs.save(PRESENTATION_PPTX)


def build_qa_docx():
    doc = Document()
    doc.add_paragraph('Q&A for Final Project Submission').bold = True
    doc.add_paragraph()

    questions = [
        ('What problem does the project address?', 'The project addresses the challenge of identifying members likely to incur high healthcare costs before those costs become large and difficult to manage.'),
        ('Why is this useful to healthcare organizations?', 'It helps organizations target care-management resources, intervene earlier, and reduce avoidable spending for high-risk members.'),
        ('What data was used?', 'The project used a realistic healthcare-style synthetic dataset with features such as age, gender, visit counts, ER visits, hospitalizations, and cost.'),
        ('Which model performed best?', 'Ridge Regression produced the strongest results, with an R² of 0.9687 and an RMSE of $1,229.62.'),
        ('What were the most important indicators of high cost?', 'Hospitalizations and emergency room visits were the strongest predictors, followed by age and utilization patterns.'),
        ('What are the limitations and next steps?', 'The dataset is synthetic and relatively small, so future work should use more diverse real-world data and explore deployment in an operational dashboard.'),
        ('How were ethical considerations addressed?', 'The project emphasized patient privacy, responsible use of synthetic versus real healthcare data, bias testing across demographic groups, fairness in decision support, and the importance of human oversight when predictive models guide healthcare interventions.')
    ]

    for q, a in questions:
        p = doc.add_paragraph()
        p.add_run(f'Q: {q}\n').bold = True
        p.add_run(f'A: {a}')
        doc.add_paragraph()

    doc.save(QA_DOCX)


def build_audio():
    text = '''Hello and welcome to this presentation on predicting high-cost healthcare members. This project focuses on a practical healthcare problem: a small group of members often drives a large share of overall spending. By identifying those individuals earlier, organizations can intervene more effectively and improve both financial and clinical outcomes. The goal of this work was to build a predictive model that estimates future healthcare cost and highlights members who may need proactive support.\n\nThe analysis used a realistic healthcare-style dataset with member-level information such as age, gender, visit counts, emergency room use, hospitalizations, and cost. The workflow followed a standard supervised learning pipeline that included data validation, exploratory analysis, feature engineering, model training, and evaluation. Several new features were created to improve the model, including age-group categories, visits-per-month, and utilization-based interaction terms.\n\nThe exploratory analysis showed that healthcare costs were highly variable and right-skewed, which is typical of spending data. A Pareto-style analysis confirmed that a relatively small share of members accounted for a large share of total costs. This supports the value of targeted intervention. The correlation analysis also showed that hospitalizations and emergency room visits had the strongest relationship with higher cost.\n\nFor modeling, several regression approaches were compared, including linear regression, ridge regression, random forest, and gradient boosting. The best-performing model was Ridge Regression, which achieved an R-squared value of 0.9687 and an RMSE of 1,229.62. These results show that the model captured the structure of the cost patterns very well and produced predictions that closely matched the observed values.\n\nThe practical value of this project is that it can support care management teams in prioritizing members who may benefit from outreach, wellness support, or early treatment planning. It also demonstrates how analytics can turn healthcare data into actionable decision support. Although the dataset used here is synthetic, the methodology is directly relevant to real-world healthcare organizations and can be extended using larger and more detailed datasets.\n\nIn summary, this project demonstrates a complete data science workflow for healthcare cost prediction. It highlights the importance of utilization-related variables, proves that strong predictive results are achievable with careful modeling, and sets up a foundation for future deployment into dashboards or operational analytics tools. Thank you.'''
    (REPORTS / 'DSC680_Presentation_Script.txt').write_text(text, encoding='utf-8')
    engine = pyttsx3.init()
    engine.setProperty('rate', 165)
    engine.save_to_file(text, str(AUDIO_WAV))
    engine.runAndWait()


def build_code_notes():
    notes = '''# Code Submission Notes

The complete project code is available in this workspace folder:

- data_collection.py
- analysis_and_visualization.py
- model_development.py
- starter_workflow.py
- config.py
- requirements.txt

The trained models are stored in the models folder and the generated outputs are in the output folder.
'''
    CODE_NOTES.write_text(notes, encoding='utf-8')


if __name__ == '__main__':
    REPORTS.mkdir(exist_ok=True)
    build_report_docx()
    build_presentation_pptx()
    build_qa_docx()
    build_audio()
    build_code_notes()
    print('Created submission assets:')
    for path in [REPORT_DOCX, PRESENTATION_PPTX, QA_DOCX, AUDIO_WAV, AUDIO_SCRIPT, CODE_NOTES]:
        print(path.name)
